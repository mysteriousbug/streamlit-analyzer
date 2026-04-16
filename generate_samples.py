"""Generate sample test files (PPTX, PDF, XLSX) matching the real-world schema.

Run once to populate sample_data/ for testing the pipeline.
"""

from pathlib import Path

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import landscape, letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

SAMPLE_DIR = Path(__file__).parent / "sample_data"


# ---------- Shared data for consistency ----------
COMMON_MONTHS = ["December'25", "January'26", "February'26"]

# DLP domain data
DLP_BREACHED = [
    ["DLP-001", "% of scoped applications without DLP monitoring",
     "Green < 3%, Amber 3% < x < 5%, Red > 5%", "Monthly",
     "4.2%", "5.8%", "6.5%",
     "Three new applications onboarded in January did not have DLP agents installed before go-live",
     "Agents deployed to 2 of 3 apps; third pending change window on 20-Feb",
     "ERR-2024-8812"],
    ["DLP-002", "% of endpoints with outdated DLP policy",
     "Green < 2%, Amber 2% < x < 4%, Red > 4%", "Monthly",
     "1.8%", "3.1%", "4.5%",
     "Policy update rollout delayed due to compatibility issue on Windows 10 LTSC fleet",
     "Rollback executed; revised policy scheduled for Feb-28 deployment",
     "M7-4421"],
]

DLP_PREV_BREACHED = [
    ["DLP-003", "% of false positives in DLP alerts",
     "Green < 10%, Amber 10% < x < 20%, Red > 20%", "Monthly",
     "22%", "18%", "12%",
     "High FP rate from overly broad regex patterns on PII detection",
     "Tuning complete for PAN and SSN patterns; continuing for email/phone",
     "ERR-2024-7103"],
]

# Governance domain data
GOV_AMBER = [
    ["GOV-101", "% of policies overdue for review",
     "Green < 5%, Amber 5% < x < 10%, Red > 10%", "Monthly",
     "3.2%", "4.1%", "7.8%",
     "Quarterly review cadence slipped due to SME unavailability during year-end",
     "Review sessions scheduled March 1-15; 12 policies in active review",
     "ERR-2024-9001"],
    ["GOV-102", "% of exceptions without approval",
     "Green < 1%, Amber 1% < x < 3%, Red > 3%", "Monthly",
     "0.5%", "1.2%", "2.1%",
     "Workflow bug caused 8 auto-approvals to skip secondary review",
     "Bug fixed Feb-10; back-filled approvals under review",
     "M7-5102"],
]


def make_pptx():
    """DLP monthly report as PPTX."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1 - Web Application Security subdomain with Breached KCI
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Web Application Security"

    # Heading above the table
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(5), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = "Breached KCI"
    p.font.bold = True
    p.font.size = Pt(16)

    _add_pptx_table(slide, DLP_BREACHED[:1], top=Inches(1.7))

    # Second table on same slide - Previous Breached KCI
    tb2 = slide.shapes.add_textbox(Inches(0.3), Inches(4.4), Inches(6), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Previous Breached KCI - monitoring"
    p2.font.bold = True
    p2.font.size = Pt(16)

    _add_pptx_table(slide, DLP_PREV_BREACHED, top=Inches(4.9))

    # Slide 2 - Endpoint Security subdomain
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Endpoint Security"
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(5), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = "Breached KCI"
    p.font.bold = True
    p.font.size = Pt(16)
    _add_pptx_table(slide, DLP_BREACHED[1:], top=Inches(1.7))

    path = SAMPLE_DIR / "DLP" / "DLP_Feb2026.pptx"
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    print(f"  wrote {path}")


def _add_pptx_table(slide, rows, top):
    """Add a table with our standard schema."""
    headers = ["ID", "Metric Name", "Threshold", "Freq",
               *COMMON_MONTHS,
               "Reason for breach", "Progress", "M7 / ERR Ref"]
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.3), top, Inches(12.7), Inches(0.5 * n_rows + 1.5)
    )
    tbl = tbl_shape.table
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            for r in para.runs:
                r.font.bold = True
                r.font.size = Pt(9)
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                for r in para.runs:
                    r.font.size = Pt(8)


def make_pdf():
    """Governance report as PDF."""
    path = SAMPLE_DIR / "Governance" / "GOV_Feb2026.pdf"
    path.parent.mkdir(parents=True, exist_ok=True)

    doc = SimpleDocTemplate(str(path), pagesize=landscape(letter))
    styles = getSampleStyleSheet()
    story = []

    story.append(Paragraph("<b>Policy Management</b>", styles["Heading1"]))
    story.append(Spacer(1, 12))
    story.append(Paragraph("<b>Amber KCI</b>", styles["Heading2"]))
    story.append(Spacer(1, 8))

    headers = ["ID", "Metric Name", "Threshold", "Freq",
               *COMMON_MONTHS,
               "Reason for breach", "Progress", "M7 / ERR Ref"]
    data = [headers] + GOV_AMBER

    tbl = Table(data, repeatRows=1, colWidths=[
        40, 140, 100, 40, 55, 55, 55, 150, 120, 55
    ])
    tbl.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#305496")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 7),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
    ]))
    story.append(tbl)

    doc.build(story)
    print(f"  wrote {path}")


def make_xlsx():
    """Third-Party Security Assessment report as XLSX."""
    wb = Workbook()
    ws = wb.active
    ws.title = "TPSA Metrics"

    # Section title
    ws["A1"] = "Third-Party Onboarding"
    ws["A1"].font = Font(bold=True, size=14)

    # Table type title
    ws["A3"] = "Breached KCI"
    ws["A3"].font = Font(bold=True, size=12)

    # Headers
    headers = ["ID", "Metric Name", "Threshold", "Freq",
               *COMMON_MONTHS,
               "Reason for breach", "Progress", "M7 / ERR Ref"]
    for col, h in enumerate(headers, start=1):
        c = ws.cell(row=4, column=col, value=h)
        c.font = Font(bold=True, color="FFFFFF")
        c.fill = PatternFill(start_color="305496", end_color="305496", fill_type="solid")

    # Data rows
    tpsa_rows = [
        ["TPSA-201", "% of high-risk vendors without completed assessment",
         "Green < 5%, Amber 5% < x < 10%, Red > 10%", "Monthly",
         "6%", "9%", "13%",
         "Backlog from Q4 onboardings combined with assessor capacity constraints",
         "Two additional assessors onboarded Feb-05; backlog reduction plan in place",
         "ERR-2025-0142"],
        ["TPSA-202", "% of vendors with expired SOC2 reports",
         "Green < 2%, Amber 2% < x < 5%, Red > 5%", "Monthly",
         "1.1%", "1.5%", "1.8%",
         "", "", "M7-6201"],
    ]
    for r_idx, row in enumerate(tpsa_rows, start=5):
        for c_idx, val in enumerate(row, start=1):
            ws.cell(row=r_idx, column=c_idx, value=val)

    # Second table on the same sheet
    ws["A9"] = "Previous Amber KRI - monitoring"
    ws["A9"].font = Font(bold=True, size=12)

    for col, h in enumerate(headers, start=1):
        c = ws.cell(row=10, column=col, value=h)
        c.font = Font(bold=True, color="FFFFFF")
        c.fill = PatternFill(start_color="305496", end_color="305496", fill_type="solid")

    prev_amber = [
        ["TPSA-203", "% of vendors not responding to annual attestation",
         "Green < 8%, Amber 8% < x < 15%, Red > 15%", "Monthly",
         "14%", "11%", "7%",
         "Improved after escalation to vendor management and contract-tied enforcement",
         "Monitoring for sustained improvement over next two cycles",
         "ERR-2024-6801"],
    ]
    for r_idx, row in enumerate(prev_amber, start=11):
        for c_idx, val in enumerate(row, start=1):
            ws.cell(row=r_idx, column=c_idx, value=val)

    # Column widths
    widths = [12, 40, 30, 10, 14, 14, 14, 40, 40, 16]
    for i, w in enumerate(widths, start=1):
        ws.column_dimensions[ws.cell(row=1, column=i).column_letter].width = w

    path = SAMPLE_DIR / "TPSA" / "TPSA_Feb2026.xlsx"
    path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(path))
    print(f"  wrote {path}")


if __name__ == "__main__":
    SAMPLE_DIR.mkdir(exist_ok=True)
    print("Generating sample files:")
    make_pptx()
    make_pdf()
    make_xlsx()
    print("Done.")
