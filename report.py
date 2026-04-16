"""Generate Excel and PPTX reports from flagged metrics."""

from __future__ import annotations

from collections import Counter, defaultdict
from datetime import datetime
from pathlib import Path

from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from pptx import Presentation
from pptx.util import Inches, Pt

from models import MetricRow, Status

# RAG colors
_GREEN_FILL = PatternFill(start_color="C6EFCE", end_color="C6EFCE", fill_type="solid")
_AMBER_FILL = PatternFill(start_color="FFEB9C", end_color="FFEB9C", fill_type="solid")
_RED_FILL = PatternFill(start_color="FFC7CE", end_color="FFC7CE", fill_type="solid")
_HEADER_FILL = PatternFill(start_color="305496", end_color="305496", fill_type="solid")
_HEADER_FONT = Font(bold=True, color="FFFFFF")


def _fill_for_status(status: Status):
    return {
        Status.GREEN: _GREEN_FILL,
        Status.AMBER: _AMBER_FILL,
        Status.RED: _RED_FILL,
    }.get(status)


def write_excel_report(metrics: list[MetricRow], output_path: Path) -> None:
    """Write a workbook with Summary / Flagged / Raw sheets."""
    wb = Workbook()

    # ---------- Summary sheet ----------
    ws = wb.active
    ws.title = "Summary"

    total = len(metrics)
    flagged = [m for m in metrics if m.is_flagged()]
    red_count = sum(1 for m in flagged if m.latest_status()[1] == Status.RED)
    amber_count = sum(1 for m in flagged if m.latest_status()[1] == Status.AMBER)
    worsening = sum(1 for m in flagged if m.trend() == "Worsening")

    ws["A1"] = "Risk Metrics Summary"
    ws["A1"].font = Font(bold=True, size=14)
    ws["A3"] = "Generated"
    ws["B3"] = datetime.now().strftime("%Y-%m-%d %H:%M")
    ws["A4"] = "Total metrics reviewed"
    ws["B4"] = total
    ws["A5"] = "Flagged (Amber + Red + Worsening)"
    ws["B5"] = len(flagged)
    ws["A6"] = "  Red breaches (latest month)"
    ws["B6"] = red_count
    ws["A7"] = "  Amber breaches (latest month)"
    ws["B7"] = amber_count
    ws["A8"] = "  Worsening trend (any severity)"
    ws["B8"] = worsening

    # Breakdown by domain / subdomain
    ws["A10"] = "Breakdown by Domain / Subdomain"
    ws["A10"].font = Font(bold=True)
    ws.append([])
    headers = ["Domain", "Subdomain", "Total", "Red", "Amber", "Flagged"]
    header_row_idx = 12
    for col, h in enumerate(headers, start=1):
        c = ws.cell(row=header_row_idx, column=col, value=h)
        c.fill = _HEADER_FILL
        c.font = _HEADER_FONT

    by_group: dict = defaultdict(lambda: {"total": 0, "red": 0, "amber": 0, "flagged": 0})
    for m in metrics:
        key = (m.domain or "(unknown)", m.subdomain or "(unknown)")
        by_group[key]["total"] += 1
        _, s = m.latest_status()
        if s == Status.RED:
            by_group[key]["red"] += 1
        elif s == Status.AMBER:
            by_group[key]["amber"] += 1
        if m.is_flagged():
            by_group[key]["flagged"] += 1

    r = header_row_idx + 1
    for (domain, sub), counts in sorted(by_group.items()):
        ws.cell(row=r, column=1, value=domain)
        ws.cell(row=r, column=2, value=sub)
        ws.cell(row=r, column=3, value=counts["total"])
        ws.cell(row=r, column=4, value=counts["red"]).fill = _RED_FILL if counts["red"] else PatternFill()
        ws.cell(row=r, column=5, value=counts["amber"]).fill = _AMBER_FILL if counts["amber"] else PatternFill()
        ws.cell(row=r, column=6, value=counts["flagged"])
        r += 1

    # Top reasons
    ws.cell(row=r + 2, column=1, value="Top Reasons for Breach").font = Font(bold=True)
    reasons = [m.reason for m in flagged if m.reason and m.reason.strip()]
    top_reasons = Counter(reasons).most_common(10)
    r += 3
    for reason, count in top_reasons:
        ws.cell(row=r, column=1, value=reason)
        ws.cell(row=r, column=2, value=count)
        r += 1

    for col_letter in ["A", "B", "C", "D", "E", "F"]:
        ws.column_dimensions[col_letter].width = 28

    # ---------- Flagged Metrics sheet ----------
    ws2 = wb.create_sheet("Flagged Metrics")
    cols = [
        "Domain", "Subdomain", "Table Type", "ID", "Metric Name",
        "Latest Month", "Latest Value", "Status", "Trend",
        "Threshold (raw)", "Reason", "Progress", "M7 / ERR Ref", "Source File",
    ]
    for col, h in enumerate(cols, start=1):
        c = ws2.cell(row=1, column=col, value=h)
        c.fill = _HEADER_FILL
        c.font = _HEADER_FONT
        c.alignment = Alignment(vertical="top", wrap_text=True)

    for i, m in enumerate(flagged, start=2):
        latest_mv, status = m.latest_status()
        ws2.cell(row=i, column=1, value=m.domain)
        ws2.cell(row=i, column=2, value=m.subdomain)
        ws2.cell(row=i, column=3, value=m.table_type.value)
        ws2.cell(row=i, column=4, value=m.metric_id)
        ws2.cell(row=i, column=5, value=m.metric_name)
        ws2.cell(row=i, column=6, value=latest_mv.label if latest_mv else "")
        ws2.cell(row=i, column=7, value=latest_mv.value if latest_mv else None)
        status_cell = ws2.cell(row=i, column=8, value=status.name)
        fill = _fill_for_status(status)
        if fill:
            status_cell.fill = fill
        ws2.cell(row=i, column=9, value=m.trend())
        ws2.cell(row=i, column=10, value=m.threshold.raw_text)
        ws2.cell(row=i, column=11, value=m.reason)
        ws2.cell(row=i, column=12, value=m.progress)
        ws2.cell(row=i, column=13, value=m.ref)
        ws2.cell(row=i, column=14, value=Path(m.source_file).name)

    widths = [18, 22, 26, 12, 40, 14, 12, 10, 12, 28, 40, 40, 18, 30]
    for col_idx, w in enumerate(widths, start=1):
        ws2.column_dimensions[ws2.cell(row=1, column=col_idx).column_letter].width = w
    ws2.freeze_panes = "A2"

    # ---------- Raw Extracted sheet ----------
    ws3 = wb.create_sheet("Raw Extracted")
    # Determine all month labels across the dataset, sorted by date
    all_months = []
    seen = set()
    for m in metrics:
        for mv in m.months:
            key = (mv.month, mv.label)
            if key not in seen:
                seen.add(key)
                all_months.append((mv.month, mv.label))
    all_months.sort(key=lambda x: x[0])
    month_labels = [lbl for _, lbl in all_months]

    raw_cols = [
        "Domain", "Subdomain", "Table Type", "ID", "Metric Name",
        "Threshold", "Frequency", *month_labels,
        "Reason", "Progress", "M7 / ERR Ref", "Source File",
    ]
    for col, h in enumerate(raw_cols, start=1):
        c = ws3.cell(row=1, column=col, value=h)
        c.fill = _HEADER_FILL
        c.font = _HEADER_FONT

    for i, m in enumerate(metrics, start=2):
        ws3.cell(row=i, column=1, value=m.domain)
        ws3.cell(row=i, column=2, value=m.subdomain)
        ws3.cell(row=i, column=3, value=m.table_type.value)
        ws3.cell(row=i, column=4, value=m.metric_id)
        ws3.cell(row=i, column=5, value=m.metric_name)
        ws3.cell(row=i, column=6, value=m.threshold.raw_text)
        ws3.cell(row=i, column=7, value=m.frequency)
        month_by_label = {mv.label: mv for mv in m.months}
        for j, lbl in enumerate(month_labels):
            col_idx = 8 + j
            mv = month_by_label.get(lbl)
            if mv is None:
                continue
            cell = ws3.cell(row=i, column=col_idx, value=mv.value)
            status = m.threshold.classify(mv.value)
            fill = _fill_for_status(status)
            if fill:
                cell.fill = fill
        ws3.cell(row=i, column=8 + len(month_labels), value=m.reason)
        ws3.cell(row=i, column=9 + len(month_labels), value=m.progress)
        ws3.cell(row=i, column=10 + len(month_labels), value=m.ref)
        ws3.cell(row=i, column=11 + len(month_labels), value=Path(m.source_file).name)

    ws3.freeze_panes = "A2"

    wb.save(str(output_path))


def write_pptx_summary(metrics: list[MetricRow], output_path: Path) -> None:
    """Write a short summary presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    flagged = [m for m in metrics if m.is_flagged()]
    red = [m for m in flagged if m.latest_status()[1] == Status.RED]
    amber = [m for m in flagged if m.latest_status()[1] == Status.AMBER]

    # Slide 1: headline
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Risk Metrics Summary"
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True

    def add_line(text, bold=False, size=18):
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold

    tf.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    tf.paragraphs[0].font.size = Pt(12)
    add_line("")
    add_line(f"Total metrics reviewed: {len(metrics)}", bold=True, size=20)
    add_line(f"Red breaches (latest month): {len(red)}", size=18)
    add_line(f"Amber breaches (latest month): {len(amber)}", size=18)
    add_line(f"Metrics with worsening trend: {sum(1 for m in flagged if m.trend() == 'Worsening')}", size=18)

    add_line("")
    add_line("Breakdown by subdomain:", bold=True, size=16)
    by_sub: Counter = Counter()
    for m in flagged:
        by_sub[m.subdomain or "(unknown)"] += 1
    for sub, count in by_sub.most_common(8):
        add_line(f"  • {sub}: {count}", size=14)

    # Slide 2: Red breaches table
    if red:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Red Breaches ({len(red)})"
        _add_breach_table(slide, red[:15])

    # Slide 3: Amber breaches table
    if amber:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Amber Breaches ({len(amber)})"
        _add_breach_table(slide, amber[:15])

    # Slide 4: Top reasons
    reasons = Counter(m.reason for m in flagged if m.reason and m.reason.strip())
    if reasons:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Top Root-Cause Themes"
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.text = ""
        for reason, count in reasons.most_common(8):
            p = tf.add_paragraph()
            p.text = f"• ({count}) {reason}"
            p.font.size = Pt(14)

    prs.save(str(output_path))


def _add_breach_table(slide, metrics: list[MetricRow]) -> None:
    """Add a formatted breach table to the slide."""
    headers = ["Subdomain", "ID", "Metric", "Latest", "Trend", "Reason"]
    n_rows = len(metrics) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.3), Inches(1.4), Inches(12.7), Inches(5.8)
    )
    tbl = tbl_shape.table
    widths = [Inches(1.8), Inches(1.0), Inches(4.0), Inches(1.2), Inches(1.2), Inches(3.5)]
    for i, w in enumerate(widths):
        tbl.columns[i].width = w

    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(11)

    for r_idx, m in enumerate(metrics, start=1):
        latest_mv, _ = m.latest_status()
        latest_str = f"{latest_mv.value:.2f}" if latest_mv and latest_mv.value is not None else ""
        row_vals = [
            m.subdomain or "",
            m.metric_id,
            (m.metric_name[:80] + "...") if len(m.metric_name) > 80 else m.metric_name,
            latest_str,
            m.trend(),
            (m.reason[:100] + "...") if len(m.reason) > 100 else (m.reason or ""),
        ]
        for c_idx, v in enumerate(row_vals):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(v)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(9)
